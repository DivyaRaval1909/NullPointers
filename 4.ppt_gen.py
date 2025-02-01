from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import textwrap

def create_ppt_with_padding(title, ai_content_dict, output_file="generated_presentation.pptx"):
    prs = Presentation()

    
    background_color = RGBColor(240, 248, 255)  # Light blue
    title_bg_color = RGBColor(0, 0, 139)  # Dark blue
    title_text_color = RGBColor(255, 255, 255)  
    content_text_color = RGBColor(50, 50, 50) 

  
    slide_layout = prs.slide_layouts[5]  
    slide = prs.slides.add_slide(slide_layout)
    

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)  
    title_frame.paragraphs[0].alignment = 1 

   
    max_lines_per_slide = 8  

    for topic, content in ai_content_dict.items():
        wrapped_content = textwrap.wrap(content, width=80)  
        total_lines = len(wrapped_content)

        first_slide_for_topic = True  
        current_slide = None
        text_frame = None

      
        for i in range(total_lines):
            if i % max_lines_per_slide == 0: 
                current_slide = prs.slides.add_slide(prs.slide_layouts[5])  

               
                current_slide.background.fill.solid()
                current_slide.background.fill.fore_color.rgb = background_color

                
                if first_slide_for_topic:
                    title_box = current_slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(1))
                    title_frame = title_box.text_frame
                    title_frame.text = topic
                    title_frame.paragraphs[0].font.size = Pt(32)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = title_text_color

                    title_box.fill.solid()
                    title_box.fill.fore_color.rgb = title_bg_color  
                    first_slide_for_topic = False  

             
                content_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Pt(10)  
                text_frame.margin_right = Pt(10) 
                text_frame.margin_top = Pt(10)
                text_frame.margin_bottom = Pt(10)  

                
            p = text_frame.add_paragraph()
            p.text = f"• {wrapped_content[i].strip()}"
            p.space_after = Pt(10)
            p.font.size = Pt(22)
            p.font.color.rgb = content_text_color 


    slide_layout = prs.slide_layouts[5]  
    slide = prs.slides.add_slide(slide_layout)

   
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color

    
    thank_you_text = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(5), Inches(1.5))
    thank_you_frame = thank_you_text.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_frame.paragraphs[0].font.size = Pt(42)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].font.color.rgb = title_bg_color  
    thank_you_frame.paragraphs[0].alignment = 1 

    
    prs.save(output_file)
    print(f"PowerPoint saved as {output_file}")




create_ppt_with_padding(title, ai_content_dict)